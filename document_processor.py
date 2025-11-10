"""
文档处理模块：负责加载和分割文本文档
支持多种文件格式：TXT, PDF, Word, Excel
"""
from langchain_text_splitters import RecursiveCharacterTextSplitter
from typing import List, Dict
import os
import mimetypes



class DocumentProcessor:
    """文档处理器，用于加载和分割文本"""
    
    def __init__(self, chunk_size: int = 500, chunk_overlap: int = 50):
        """
        初始化文档处理器
        
        Args:
            chunk_size: 文本块大小
            chunk_overlap: 文本块重叠大小
        """
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", "。", "！", "？", "；", " ", ""]
        )
    
    def _get_file_type(self, file_path: str) -> str:
        """
        获取文件类型
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件类型（txt, pdf, docx, xlsx等）
        """
        _, ext = os.path.splitext(file_path.lower())
        return ext[1:] if ext else ''
    
    def load_text_file(self, file_path: str) -> str:
        """
        加载文本文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件内容
        """
        try:
            # 尝试多种编码
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise Exception("无法解码文件，请检查文件编码")
        except Exception as e:
            raise Exception(f"加载文本文件失败: {str(e)}")
    
    def load_pdf_file(self, file_path: str) -> str:
        """
        加载PDF文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件内容
        """
        try:
            import pdfplumber
            
            text_content = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
            
            return '\n\n'.join(text_content)
        except ImportError:
            # 如果pdfplumber不可用，尝试使用pypdf
            try:
                from pypdf import PdfReader
                text_content = []
                with open(file_path, 'rb') as f:
                    pdf_reader = PdfReader(f)
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            text_content.append(text)
                return '\n\n'.join(text_content)
            except ImportError:
                raise Exception("PDF处理库未安装，请安装pdfplumber或pypdf")
        except Exception as e:
            raise Exception(f"加载PDF文件失败: {str(e)}")
    
    def load_word_file(self, file_path: str) -> str:
        """
        加载Word文件（.docx）
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件内容
        """
        try:
            from docx import Document
            
            doc = Document(file_path)
            text_content = []
            
            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(' | '.join(row_text))
            
            return '\n\n'.join(text_content)
        except ImportError:
            raise Exception("python-docx库未安装，请安装: pip install python-docx")
        except Exception as e:
            raise Exception(f"加载Word文件失败: {str(e)}")
    
    def load_ppt_file(self, file_path: str) -> str:
        """
        加载PowerPoint文件（.pptx）
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件内容
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            # 遍历所有幻灯片
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_texts = []
                slide_texts.append(f"=== 幻灯片 {slide_idx} ===")
                
                # 提取幻灯片标题
                try:
                    if slide.shapes.title:
                        title = slide.shapes.title.text.strip()
                        if title:
                            slide_texts.append(f"标题: {title}")
                except:
                    pass  # 某些幻灯片可能没有标题
                
                # 提取所有形状中的文本
                for shape in slide.shapes:
                    try:
                        # 跳过标题（已经处理过）
                        if slide.shapes.title and shape == slide.shapes.title:
                            continue
                        
                        # 处理文本框和自动形状
                        if hasattr(shape, "text"):
                            text = shape.text.strip()
                            if text:
                                slide_texts.append(text)
                        
                        # 处理表格
                        if hasattr(shape, "has_table") and shape.has_table:
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_text = []
                                for cell in row.cells:
                                    cell_text = cell.text.strip()
                                    if cell_text:
                                        row_text.append(cell_text)
                                if row_text:
                                    table_data.append(" | ".join(row_text))
                            if table_data:
                                slide_texts.append("表格:")
                                slide_texts.extend(table_data)
                        
                        # 处理占位符（Placeholder）
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                            if hasattr(shape, "text"):
                                placeholder_text = shape.text.strip()
                                if placeholder_text:
                                    slide_texts.append(placeholder_text)
                    except Exception as e:
                        # 跳过无法处理的形状
                        continue
                
                # 提取备注
                try:
                    if slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        if notes_slide.notes_text_frame:
                            notes = notes_slide.notes_text_frame.text.strip()
                            if notes:
                                slide_texts.append(f"备注: {notes}")
                except:
                    pass  # 某些幻灯片可能没有备注
                
                # 如果有内容，添加到总内容中
                if len(slide_texts) > 1:  # 至少有幻灯片编号
                    text_content.extend(slide_texts)
                    text_content.append("")  # 添加空行分隔
            
            if not text_content:
                raise Exception("PowerPoint文件中没有找到可提取的文本内容")
            
            return '\n\n'.join(text_content)
        except ImportError:
            raise Exception("python-pptx库未安装，请安装: pip install python-pptx")
        except Exception as e:
            raise Exception(f"加载PowerPoint文件失败: {str(e)}")
    
    def load_excel_file(self, file_path: str) -> str:
        """
        加载Excel文件（.xlsx, .xls）
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件内容
        """
        try:
            import pandas as pd
            
            # 读取所有工作表
            # pandas会自动根据文件扩展名选择合适的引擎
            # .xlsx使用openpyxl, .xls使用xlrd
            try:
                excel_file = pd.ExcelFile(file_path)
            except Exception as e:
                # 如果读取失败，可能是缺少xlrd库（用于.xls文件）
                if file_path.endswith('.xls'):
                    raise Exception(f"读取.xls文件失败，可能需要安装xlrd库: pip install xlrd。错误: {str(e)}")
                raise e
            
            text_content = []
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    
                    # 添加工作表名称
                    text_content.append(f"工作表: {sheet_name}")
                    
                    # 将DataFrame转换为文本
                    # 使用to_string()或者转换为CSV格式的文本
                    # 处理NaN值
                    df = df.fillna('')  # 将NaN替换为空字符串
                    text_content.append(df.to_string(index=False))
                    text_content.append("")  # 添加空行分隔
                except Exception as e:
                    # 如果某个工作表读取失败，继续处理其他工作表
                    text_content.append(f"工作表: {sheet_name} (读取失败: {str(e)})")
                    text_content.append("")
            
            return '\n\n'.join(text_content)
        except ImportError as e:
            raise Exception(f"Excel处理库未安装: {str(e)}。请安装: pip install pandas openpyxl xlrd")
        except Exception as e:
            raise Exception(f"加载Excel文件失败: {str(e)}")
    
    def load_document(self, file_path: str) -> str:
        """
        根据文件类型加载文档
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件内容
        """
        file_type = self._get_file_type(file_path)
        
        if file_type == 'txt' or file_type == '':
            return self.load_text_file(file_path)
        elif file_type == 'pdf':
            return self.load_pdf_file(file_path)
        elif file_type in ['docx', 'doc']:
            if file_type == 'doc':
                raise Exception("仅支持.docx格式，不支持旧的.doc格式。请将文件转换为.docx格式。")
            return self.load_word_file(file_path)
        elif file_type in ['xlsx', 'xls']:
            return self.load_excel_file(file_path)
        elif file_type in ['pptx']:
            return self.load_ppt_file(file_path)
        else:
            raise Exception(f"不支持的文件格式: {file_type}。支持格式: txt, pdf, docx, xlsx, xls, pptx")
    
    def split_text(self, text: str) -> List[str]:
        """
        分割文本为块
        
        Args:
            text: 原始文本
            
        Returns:
            文本块列表
        """
        chunks = self.text_splitter.split_text(text)
        return chunks
    
    def process_document(self, file_path: str) -> List[Dict[str, str]]:
        """
        处理文档：加载并分割
        
        Args:
            file_path: 文件路径
            
        Returns:
            文档块列表，每个块包含内容和元数据
        """
        # 使用新的load_document方法
        text = self.load_document(file_path)
        chunks = self.split_text(text)
        
        file_type = self._get_file_type(file_path)
        file_name = os.path.basename(file_path)
        
        documents = []
        for i, chunk in enumerate(chunks):
            documents.append({
                "content": chunk,
                "metadata": {
                    "source": file_path,
                    "chunk_index": i,
                    "file_name": file_name,
                    "file_type": file_type
                }
            })
        
        return documents
    
    def process_text(self, text: str, source: str = "input") -> List[Dict[str, str]]:
        """
        处理文本：直接分割文本
        
        Args:
            text: 原始文本
            source: 来源标识
            
        Returns:
            文档块列表
        """
        chunks = self.split_text(text)
        
        documents = []
        for i, chunk in enumerate(chunks):
            documents.append({
                "content": chunk,
                "metadata": {
                    "source": source,
                    "chunk_index": i
                }
            })
        
        return documents

